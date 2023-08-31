import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches

def select_images():
    image_files = filedialog.askopenfilenames(filetypes=[("Image files", "*.png *.jpg *.jpeg")])
    image_listbox.delete(0, tk.END)
    for image in image_files:
        image_listbox.insert(tk.END, image)

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    selected_images = image_listbox.get(0, tk.END)
    
    for image_file in selected_images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        left = (prs.slide_width - Inches(16)) / 2  # 居中图片
        top = (prs.slide_height - Inches(9)) / 2
        pic = slide.shapes.add_picture(image_file, left, top, width=Inches(16), height=Inches(9))
        
        # 调整图片大小以适应幻灯片，但不超出底图
        pic.width = min(pic.width, prs.slide_width)
        pic.height = min(pic.height, prs.slide_height)
        
    prs.save("output.pptx")
    result_label.config(text="PPTX created successfully!")

# 创建主窗口
root = tk.Tk()
root.title("Image to PPTX Converter")

# 添加选择图片按钮
select_button = tk.Button(root, text="Select Images", command=select_images)
select_button.pack(pady=10)

# 显示已选择的图片列表
image_listbox = tk.Listbox(root, selectmode=tk.MULTIPLE)
image_listbox.pack()

# 添加创建PPTX按钮
create_button = tk.Button(root, text="Create PPTX", command=create_pptx)
create_button.pack(pady=10)

# 显示操作结果
result_label = tk.Label(root, text="")
result_label.pack()

# 运行主循环
root.mainloop()
