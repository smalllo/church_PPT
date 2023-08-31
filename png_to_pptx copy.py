import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches

def select_images():
    image_files = filedialog.askopenfilenames(filetypes=[("Image files", "*.png *.jpg *.jpeg")])
    image_listbox.delete(0, tk.END)
    for image in image_files:
        image_listbox.insert(tk.END, image)

def select_output_path():
    output_file = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PPTX files", "*.pptx")])
    if output_file:
        output_path_entry.delete(0, tk.END)
        output_path_entry.insert(0, output_file)

def create_pptx():
    # 获取输出路径
    output_file = output_path_entry.get()

    if not output_file:
        result_label.config(text="請選擇輸出PPTX檔案位置")
        return

    prs = Presentation()
    
    # 選取ppt比例
    selected_ratio = ratio_var.get()
    if selected_ratio == "4:3":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        selected_images = image_listbox.get(0, tk.END)
        
        for image_file in selected_images:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            left = (prs.slide_width - Inches(10)) / 2
            top = (prs.slide_height - Inches(7.5)) / 2
            pic = slide.shapes.add_picture(image_file, left, top, width=Inches(10), height=Inches(7.5))
            pic.width = min(pic.width, prs.slide_width)
            pic.height = min(pic.height, prs.slide_height)

    else:
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

        selected_images = image_listbox.get(0, tk.END)
        
        for image_file in selected_images:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            left = (prs.slide_width - Inches(16)) / 2
            top = (prs.slide_height - Inches(9)) / 2
            pic = slide.shapes.add_picture(image_file, left, top, width=Inches(16), height=Inches(9))
            pic.width = min(pic.width, prs.slide_width)
            pic.height = min(pic.height, prs.slide_height)
            
    prs.save(output_file)
    result_label.config(text="PPTX創建成功！")

root = tk.Tk()
root.title("png轉pptx By小羅")
root.geometry("600x400")

select_button = tk.Button(root, text="選擇圖片", command=select_images)
select_button.pack()

hello_label = tk.Label(root, text="以選取圖片")
hello_label.pack()

image_listbox = tk.Listbox(root, selectmode=tk.MULTIPLE)
image_listbox.pack()

hello_label = tk.Label(root, text="選擇投影片比例")
hello_label.pack()

ratio_var = tk.StringVar()
ratio_frame = tk.Frame(root)
ratio_frame.pack()
four_three_radio = tk.Radiobutton(ratio_frame, text="4:3", variable=ratio_var, value="4:3")
sixteen_nine_radio = tk.Radiobutton(ratio_frame, text="16:9", variable=ratio_var, value="16:9")

four_three_radio.pack(side=tk.LEFT)
sixteen_nine_radio.pack(side=tk.LEFT)

output_path_frame = tk.Frame(root)
output_path_frame.pack(pady=10)

output_path_label = tk.Label(output_path_frame, text="選擇輸出PPTX檔案位置")
output_path_label.pack()

output_path_entry = tk.Entry(output_path_frame, width=40)
output_path_entry.pack(side=tk.LEFT)

output_path_button = tk.Button(output_path_frame, text="瀏覽", command=select_output_path)
output_path_button.pack(side=tk.LEFT)

create_button = tk.Button(root, text="開始", command=create_pptx)
create_button.pack(pady=10)

result_label = tk.Label(root, text="")
result_label.pack()

root.mainloop()
